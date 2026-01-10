"""Tests for document conversion handlers."""

import os
import tempfile
import pytest
from io import BytesIO
from unittest.mock import MagicMock, patch, PropertyMock

# Import conversion libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib import colors
from pypdf import PdfReader


def extract_pdf_text(pdf_bytes):
    """Extract text from PDF bytes using pypdf."""
    reader = PdfReader(BytesIO(pdf_bytes))
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text


def create_mock_handler():
    """Create a handler instance with mocked log property."""
    from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler
    handler = object.__new__(DocumentConverterHandler)
    # Set the _log attribute directly to bypass property
    object.__setattr__(handler, '_log', MagicMock())
    return handler


class TestPPTXConversion:
    """Test PPTX to PDF conversion functionality."""

    @pytest.fixture
    def sample_pptx(self):
        """Create a sample PPTX file for testing."""
        prs = Presentation()

        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Test Presentation"
        subtitle.text = "Created for unit testing"

        # Add content slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide2.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Test Slide"
        tf = body_shape.text_frame
        tf.text = "First bullet point"
        p = tf.add_paragraph()
        p.text = "Second bullet point"
        p.level = 1
        p = tf.add_paragraph()
        p.text = "Polish characters: ąćęłńóśźż"
        p.level = 0

        # Save to temp file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            yield f.name

        # Cleanup
        os.unlink(f.name)

    def test_pptx_loads_correctly(self, sample_pptx):
        """Test that sample PPTX file loads correctly."""
        prs = Presentation(sample_pptx)
        assert len(prs.slides) == 2

    def test_pptx_has_slide_dimensions(self, sample_pptx):
        """Test that PPTX has valid slide dimensions."""
        prs = Presentation(sample_pptx)
        assert prs.slide_width > 0
        assert prs.slide_height > 0

    def test_pptx_conversion_produces_pdf(self, sample_pptx):
        """Test that PPTX conversion produces valid PDF bytes."""
        from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler

        with patch.object(DocumentConverterHandler, 'log', new_callable=PropertyMock) as mock_log:
            mock_log.return_value = MagicMock()
            handler = object.__new__(DocumentConverterHandler)
            pdf_bytes = handler._convert_pptx_to_pdf(sample_pptx)

        # Check that we got PDF bytes
        assert pdf_bytes is not None
        assert len(pdf_bytes) > 0
        # PDF files start with %PDF
        assert pdf_bytes[:4] == b'%PDF'

    def test_pptx_conversion_multiple_slides(self, sample_pptx):
        """Test that conversion handles multiple slides."""
        from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler

        with patch.object(DocumentConverterHandler, 'log', new_callable=PropertyMock) as mock_log:
            mock_log.return_value = MagicMock()
            handler = object.__new__(DocumentConverterHandler)
            pdf_bytes = handler._convert_pptx_to_pdf(sample_pptx)

        # PDF should contain multiple pages
        assert pdf_bytes is not None
        # Check PDF contains page references
        assert b'/Page' in pdf_bytes


class TestDOCXConversion:
    """Test DOCX to PDF conversion functionality."""

    @pytest.fixture
    def sample_docx(self):
        """Create a sample DOCX file for testing."""
        doc = Document()
        doc.add_heading('Test Document', 0)
        doc.add_paragraph('This is a test paragraph.')
        doc.add_paragraph('Polish characters: ąćęłńóśźż')
        doc.add_heading('Section 1', level=1)
        doc.add_paragraph('More content here.')

        # Save to temp file
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
            doc.save(f.name)
            yield f.name

        # Cleanup
        os.unlink(f.name)

    def test_docx_loads_correctly(self, sample_docx):
        """Test that sample DOCX file loads correctly."""
        doc = Document(sample_docx)
        assert len(doc.paragraphs) > 0

    def test_docx_conversion_produces_pdf(self, sample_docx):
        """Test that DOCX conversion produces valid PDF bytes."""
        from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler

        with patch.object(DocumentConverterHandler, 'log', new_callable=PropertyMock) as mock_log:
            mock_log.return_value = MagicMock()
            handler = object.__new__(DocumentConverterHandler)
            pdf_bytes = handler._convert_docx_to_pdf(sample_docx)

        # Check that we got PDF bytes
        assert pdf_bytes is not None
        assert len(pdf_bytes) > 0
        # PDF files start with %PDF
        assert pdf_bytes[:4] == b'%PDF'


class TestFileTypeRouting:
    """Test file type detection and routing."""

    def test_unsupported_ppt_format(self):
        """Test that legacy PPT format raises appropriate error."""
        from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler

        with patch.object(DocumentConverterHandler, 'log', new_callable=PropertyMock) as mock_log:
            mock_log.return_value = MagicMock()
            handler = object.__new__(DocumentConverterHandler)
            with pytest.raises(Exception) as exc_info:
                handler._convert_to_pdf('/path/to/file.ppt')
            assert 'PPT format not supported' in str(exc_info.value)

    def test_unsupported_doc_format(self):
        """Test that legacy DOC format raises appropriate error."""
        from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler

        with patch.object(DocumentConverterHandler, 'log', new_callable=PropertyMock) as mock_log:
            mock_log.return_value = MagicMock()
            handler = object.__new__(DocumentConverterHandler)
            with pytest.raises(Exception) as exc_info:
                handler._convert_to_pdf('/path/to/file.doc')
            assert 'DOC' in str(exc_info.value) or 'not supported' in str(exc_info.value)

    def test_unsupported_rtf_format(self):
        """Test that RTF format raises appropriate error."""
        from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler

        with patch.object(DocumentConverterHandler, 'log', new_callable=PropertyMock) as mock_log:
            mock_log.return_value = MagicMock()
            handler = object.__new__(DocumentConverterHandler)
            with pytest.raises(Exception) as exc_info:
                handler._convert_to_pdf('/path/to/file.rtf')
            assert 'RTF' in str(exc_info.value) or 'not supported' in str(exc_info.value)

    def test_unknown_file_type(self):
        """Test that unknown file types raise appropriate error."""
        from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler

        with patch.object(DocumentConverterHandler, 'log', new_callable=PropertyMock) as mock_log:
            mock_log.return_value = MagicMock()
            handler = object.__new__(DocumentConverterHandler)
            with pytest.raises(Exception) as exc_info:
                handler._convert_to_pdf('/path/to/file.xyz')
            assert 'Unsupported file type' in str(exc_info.value)


class TestUnicodeFonts:
    """Test Unicode font registration."""

    def test_font_registration_method_exists(self):
        """Test that font registration method exists."""
        from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler

        # Method should exist and be callable
        assert hasattr(DocumentConverterHandler, '_register_unicode_fonts')
        assert callable(DocumentConverterHandler._register_unicode_fonts)

    def test_font_registration_runs_without_error(self):
        """Test that font registration runs without error."""
        from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler

        with patch.object(DocumentConverterHandler, 'log', new_callable=PropertyMock) as mock_log:
            mock_log.return_value = MagicMock()
            handler = object.__new__(DocumentConverterHandler)
            # Should not raise
            handler._register_unicode_fonts()


class TestImportAvailability:
    """Test that required imports are available."""

    def test_reportlab_available(self):
        """Test that reportlab imports are available."""
        from jupyterlab_doc_reader_extension.handlers import REPORTLAB_AVAILABLE
        assert REPORTLAB_AVAILABLE is True

    def test_docx_available(self):
        """Test that python-docx imports are available."""
        from jupyterlab_doc_reader_extension.handlers import DOCX_AVAILABLE
        assert DOCX_AVAILABLE is True

    def test_pptx_available(self):
        """Test that python-pptx imports are available."""
        from jupyterlab_doc_reader_extension.handlers import PPTX_AVAILABLE
        assert PPTX_AVAILABLE is True

    def test_pil_available(self):
        """Test that PIL imports are available."""
        from jupyterlab_doc_reader_extension.handlers import PIL_AVAILABLE
        assert PIL_AVAILABLE is True


class TestPDFGeneration:
    """Test basic PDF generation with reportlab."""

    def test_create_simple_pdf(self):
        """Test that we can create a simple PDF."""
        pdf_buffer = BytesIO()
        c = canvas.Canvas(pdf_buffer)
        c.drawString(100, 750, "Test PDF")
        c.showPage()
        c.save()

        pdf_bytes = pdf_buffer.getvalue()
        assert pdf_bytes[:4] == b'%PDF'
        assert len(pdf_bytes) > 100

    def test_pdf_with_colors(self):
        """Test PDF generation with colors."""
        pdf_buffer = BytesIO()
        c = canvas.Canvas(pdf_buffer)
        c.setFillColor(colors.red)
        c.drawString(100, 750, "Red text")
        c.setFillColor(colors.blue)
        c.drawString(100, 700, "Blue text")
        c.showPage()
        c.save()

        pdf_bytes = pdf_buffer.getvalue()
        assert pdf_bytes[:4] == b'%PDF'


class TestDocumentOrderPreservation:
    """Test that document elements are rendered in correct order."""

    @pytest.fixture
    def docx_with_inline_table(self):
        """Create a DOCX with paragraph, table, paragraph (in that order)."""
        doc = Document()
        doc.add_heading('Title', level=1)
        doc.add_paragraph('Paragraph before table.')

        # Add a table
        table = doc.add_table(rows=2, cols=3)
        table.cell(0, 0).text = 'Header1'
        table.cell(0, 1).text = 'Header2'
        table.cell(0, 2).text = 'Header3'
        table.cell(1, 0).text = 'Data1'
        table.cell(1, 1).text = 'Data2'
        table.cell(1, 2).text = 'Data3'

        doc.add_paragraph('Paragraph after table.')
        doc.add_heading('Section 2', level=2)
        doc.add_paragraph('Final paragraph.')

        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
            doc.save(f.name)
            yield f.name

        os.unlink(f.name)

    def test_table_position_preserved(self, docx_with_inline_table):
        """Test that table appears inline between paragraphs, not at end."""
        from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler

        with patch.object(DocumentConverterHandler, 'log', new_callable=PropertyMock) as mock_log:
            mock_log.return_value = MagicMock()
            handler = object.__new__(DocumentConverterHandler)
            pdf_bytes = handler._convert_docx_to_pdf(docx_with_inline_table)

        # PDF should be valid
        assert pdf_bytes[:4] == b'%PDF'

        # Extract text from PDF using pypdf
        pdf_text = extract_pdf_text(pdf_bytes)

        # Find positions of key elements
        pos_before = pdf_text.find('Paragraph before table')
        pos_header = pdf_text.find('Header1')
        pos_after = pdf_text.find('Paragraph after table')

        # All elements should be found
        assert pos_before != -1, f"'Paragraph before table' not found in PDF. Text: {pdf_text[:500]}"
        assert pos_header != -1, f"'Header1' not found in PDF. Text: {pdf_text[:500]}"
        assert pos_after != -1, f"'Paragraph after table' not found in PDF. Text: {pdf_text[:500]}"

        # Table content should appear between the two paragraphs
        assert pos_before < pos_header < pos_after, \
            "Table should appear between 'Paragraph before table' and 'Paragraph after table'"


class TestBoldItalicHandling:
    """Test bold and italic text formatting."""

    @pytest.fixture
    def docx_with_formatting(self):
        """Create a DOCX with bold and italic text."""
        doc = Document()
        doc.add_heading('Formatted Document', level=1)

        # Add paragraph with mixed formatting
        para = doc.add_paragraph()
        para.add_run('Normal text, ')
        bold_run = para.add_run('bold text')
        bold_run.bold = True
        para.add_run(', ')
        italic_run = para.add_run('italic text')
        italic_run.italic = True
        para.add_run(', and ')
        bold_italic_run = para.add_run('bold italic')
        bold_italic_run.bold = True
        bold_italic_run.italic = True
        para.add_run('.')

        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
            doc.save(f.name)
            yield f.name

        os.unlink(f.name)

    def test_formatted_text_renders(self, docx_with_formatting):
        """Test that bold and italic text is processed."""
        from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler

        with patch.object(DocumentConverterHandler, 'log', new_callable=PropertyMock) as mock_log:
            mock_log.return_value = MagicMock()
            handler = object.__new__(DocumentConverterHandler)
            pdf_bytes = handler._convert_docx_to_pdf(docx_with_formatting)

        # PDF should be valid
        assert pdf_bytes[:4] == b'%PDF'
        assert len(pdf_bytes) > 1000

        # Extract text from PDF using pypdf
        pdf_text = extract_pdf_text(pdf_bytes)

        # Check that content is present
        assert 'Normal text' in pdf_text, f"'Normal text' not found. PDF text: {pdf_text[:500]}"
        assert 'bold text' in pdf_text, f"'bold text' not found. PDF text: {pdf_text[:500]}"
        assert 'italic text' in pdf_text, f"'italic text' not found. PDF text: {pdf_text[:500]}"


class TestListHandling:
    """Test bullet and numbered list handling."""

    @pytest.fixture
    def docx_with_lists(self):
        """Create a DOCX with bullet and numbered lists."""
        doc = Document()
        doc.add_heading('Document with Lists', level=1)

        doc.add_paragraph('Introduction paragraph.')

        # Add bullet list
        doc.add_paragraph('Bullet item 1', style='List Bullet')
        doc.add_paragraph('Bullet item 2', style='List Bullet')
        doc.add_paragraph('Bullet item 3', style='List Bullet')

        doc.add_paragraph('Middle paragraph.')

        # Add numbered list
        doc.add_paragraph('Number item 1', style='List Number')
        doc.add_paragraph('Number item 2', style='List Number')
        doc.add_paragraph('Number item 3', style='List Number')

        doc.add_paragraph('Final paragraph.')

        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
            doc.save(f.name)
            yield f.name

        os.unlink(f.name)

    def test_bullet_list_renders(self, docx_with_lists):
        """Test that bullet lists are rendered with bullet characters."""
        from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler

        with patch.object(DocumentConverterHandler, 'log', new_callable=PropertyMock) as mock_log:
            mock_log.return_value = MagicMock()
            handler = object.__new__(DocumentConverterHandler)
            pdf_bytes = handler._convert_docx_to_pdf(docx_with_lists)

        # PDF should be valid
        assert pdf_bytes[:4] == b'%PDF'

        # Extract text from PDF using pypdf
        pdf_text = extract_pdf_text(pdf_bytes)

        # Check that list items are present
        assert 'Bullet item 1' in pdf_text, f"'Bullet item 1' not found. PDF text: {pdf_text[:500]}"
        assert 'Number item 1' in pdf_text, f"'Number item 1' not found. PDF text: {pdf_text[:500]}"

    def test_numbered_list_has_numbers(self, docx_with_lists):
        """Test that numbered lists have sequential numbers."""
        from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler

        with patch.object(DocumentConverterHandler, 'log', new_callable=PropertyMock) as mock_log:
            mock_log.return_value = MagicMock()
            handler = object.__new__(DocumentConverterHandler)
            pdf_bytes = handler._convert_docx_to_pdf(docx_with_lists)

        # Extract text from PDF using pypdf
        pdf_text = extract_pdf_text(pdf_bytes)

        # Check for numbered list markers (1., 2., 3.)
        assert '1.' in pdf_text and 'Number item 1' in pdf_text, f"Numbered list item 1 not found. Text: {pdf_text[:500]}"
        assert '2.' in pdf_text and 'Number item 2' in pdf_text, f"Numbered list item 2 not found. Text: {pdf_text[:500]}"
        assert '3.' in pdf_text and 'Number item 3' in pdf_text, f"Numbered list item 3 not found. Text: {pdf_text[:500]}"


class TestHeadingStyles:
    """Test heading style detection and rendering."""

    @pytest.fixture
    def docx_with_headings(self):
        """Create a DOCX with multiple heading levels."""
        doc = Document()
        doc.add_heading('Main Title', level=0)
        doc.add_paragraph('Introduction.')
        doc.add_heading('Section 1', level=1)
        doc.add_paragraph('Section 1 content.')
        doc.add_heading('Subsection 1.1', level=2)
        doc.add_paragraph('Subsection content.')
        doc.add_heading('Sub-subsection', level=3)
        doc.add_paragraph('Deep content.')

        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
            doc.save(f.name)
            yield f.name

        os.unlink(f.name)

    def test_headings_are_rendered(self, docx_with_headings):
        """Test that all heading levels are rendered."""
        from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler

        with patch.object(DocumentConverterHandler, 'log', new_callable=PropertyMock) as mock_log:
            mock_log.return_value = MagicMock()
            handler = object.__new__(DocumentConverterHandler)
            pdf_bytes = handler._convert_docx_to_pdf(docx_with_headings)

        # PDF should be valid
        assert pdf_bytes[:4] == b'%PDF'

        # Extract text from PDF using pypdf
        pdf_text = extract_pdf_text(pdf_bytes)

        # Check that all headings are present
        assert 'Main Title' in pdf_text, f"'Main Title' not found. PDF text: {pdf_text[:500]}"
        assert 'Section 1' in pdf_text, f"'Section 1' not found. PDF text: {pdf_text[:500]}"
        assert 'Subsection 1.1' in pdf_text, f"'Subsection 1.1' not found. PDF text: {pdf_text[:500]}"


class TestTableStyling:
    """Test table styling and formatting."""

    @pytest.fixture
    def docx_with_styled_table(self):
        """Create a DOCX with a styled table."""
        doc = Document()
        doc.add_heading('Table Test', level=1)

        table = doc.add_table(rows=3, cols=4)
        # Header row
        table.cell(0, 0).text = 'Field'
        table.cell(0, 1).text = 'Type'
        table.cell(0, 2).text = 'Description'
        table.cell(0, 3).text = 'Example'
        # Data rows
        table.cell(1, 0).text = 'name'
        table.cell(1, 1).text = 'string'
        table.cell(1, 2).text = 'User name'
        table.cell(1, 3).text = 'John'
        table.cell(2, 0).text = 'age'
        table.cell(2, 1).text = 'int'
        table.cell(2, 2).text = 'Age in years'
        table.cell(2, 3).text = '25'

        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
            doc.save(f.name)
            yield f.name

        os.unlink(f.name)

    def test_table_content_preserved(self, docx_with_styled_table):
        """Test that table content is preserved in PDF."""
        from jupyterlab_doc_reader_extension.handlers import DocumentConverterHandler

        with patch.object(DocumentConverterHandler, 'log', new_callable=PropertyMock) as mock_log:
            mock_log.return_value = MagicMock()
            handler = object.__new__(DocumentConverterHandler)
            pdf_bytes = handler._convert_docx_to_pdf(docx_with_styled_table)

        # PDF should be valid
        assert pdf_bytes[:4] == b'%PDF'

        # Extract text from PDF using pypdf
        pdf_text = extract_pdf_text(pdf_bytes)

        # Check that table content is present
        assert 'Field' in pdf_text, f"'Field' not found. PDF text: {pdf_text[:500]}"
        assert 'Type' in pdf_text, f"'Type' not found. PDF text: {pdf_text[:500]}"
        assert 'Description' in pdf_text, f"'Description' not found. PDF text: {pdf_text[:500]}"
        assert 'name' in pdf_text, f"'name' not found. PDF text: {pdf_text[:500]}"
        assert 'string' in pdf_text, f"'string' not found. PDF text: {pdf_text[:500]}"
