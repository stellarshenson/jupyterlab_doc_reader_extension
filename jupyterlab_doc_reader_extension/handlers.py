import json
import os
import tempfile
import base64
from pathlib import Path
from io import BytesIO

from jupyter_server.base.handlers import APIHandler
from jupyter_server.utils import url_path_join
import tornado

# Reportlab imports (required for both DOCX and PPTX conversion)
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle, Image as RLImage
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
    from reportlab.lib import colors
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.pdfgen import canvas
    REPORTLAB_AVAILABLE = True
except ImportError as e:
    REPORTLAB_AVAILABLE = False
    _reportlab_import_error = str(e)

# python-docx for DOCX conversion
try:
    from docx import Document
    from html import escape as html_escape
    DOCX_AVAILABLE = True
except ImportError as e:
    Document = None
    DOCX_AVAILABLE = False
    _docx_import_error = str(e)

# python-pptx for PPTX conversion
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu as PptxEmu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError as e:
    Presentation = None
    PPTX_AVAILABLE = False
    _pptx_import_error = str(e)

# Pillow for image handling
try:
    from PIL import Image as PILImage
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False


class DocumentConverterHandler(APIHandler):
    """Handler for converting documents (DOCX, DOC, RTF, PPTX, PPT) to PDF"""

    @tornado.web.authenticated
    def post(self):
        """
        Convert a document file to PDF and return as base64-encoded data.
        Expects JSON payload: {"path": "/path/to/document.docx"}
        """
        try:
            data = self.get_json_body()
            file_path = data.get('path')

            self.log.info(f"Converting document: {file_path}")
            self.log.debug(f"Request data: {data}")

            if not file_path:
                self.set_status(400)
                self.finish(json.dumps({"error": "No file path provided"}))
                return

            # Get the full path to the file
            contents_manager = self.settings.get('contents_manager')
            if contents_manager:
                root_dir = contents_manager.root_dir
            else:
                root_dir = os.getcwd()
                self.log.warning(f"Contents manager not found, using cwd: {root_dir}")

            full_path = os.path.join(root_dir, file_path.lstrip('/'))
            self.log.debug(f"Full path: {full_path}")

            if not os.path.exists(full_path):
                self.set_status(404)
                self.finish(json.dumps({"error": f"File not found: {file_path}"}))
                return

            # Check file extension
            ext = Path(full_path).suffix.lower()
            if ext not in ['.docx', '.doc', '.rtf', '.pptx', '.ppt']:
                self.set_status(400)
                self.finish(json.dumps({"error": f"Unsupported file type: {ext}"}))
                return

            # Convert to PDF
            try:
                pdf_data = self._convert_to_pdf(full_path)
                self.log.debug(f"PDF size: {len(pdf_data)} bytes")

                # Encode as base64 for transmission
                pdf_base64 = base64.b64encode(pdf_data).decode('utf-8')

                response_data = {
                    "success": True,
                    "pdf_data": pdf_base64,
                    "filename": Path(file_path).stem + ".pdf"
                }

                self.finish(json.dumps(response_data))
                self.log.info(f"Conversion successful")

            except Exception as convert_error:
                import traceback
                error_traceback = traceback.format_exc()
                self.log.error(f"Conversion error: {str(convert_error)}\n{error_traceback}")
                self.set_status(500)
                self.finish(json.dumps({
                    "success": False,
                    "error": f"Conversion failed: {str(convert_error)}",
                    "error_type": type(convert_error).__name__,
                    "traceback": error_traceback,
                    "file_path": file_path,
                    "full_path": full_path
                }))

        except Exception as e:
            import traceback
            error_traceback = traceback.format_exc()
            self.log.error(f"Handler error: {str(e)}\n{error_traceback}")
            self.set_status(500)
            self.finish(json.dumps({
                "success": False,
                "error": str(e),
                "error_type": type(e).__name__,
                "traceback": error_traceback
            }))

    def _convert_to_pdf(self, input_path: str) -> bytes:
        """
        Convert document to PDF using python-docx/python-pptx + reportlab.
        Pure Python solution with no external system dependencies.
        Returns PDF data as bytes.
        """
        ext = Path(input_path).suffix.lower()

        # Route to appropriate converter based on file type
        if ext == '.pptx':
            return self._convert_pptx_to_pdf(input_path)
        elif ext == '.ppt':
            raise Exception(
                "Legacy PPT format not supported. "
                "Please convert to PPTX format for best results."
            )
        elif ext == '.docx':
            return self._convert_docx_to_pdf(input_path)
        elif ext in ['.doc', '.rtf']:
            raise Exception(
                f"Legacy {ext.upper()} format not supported. "
                f"Please convert to DOCX format for best results."
            )
        else:
            raise Exception(f"Unsupported file type: {ext}")

    def _convert_docx_to_pdf(self, input_path: str) -> bytes:
        """
        Convert DOCX document to PDF using python-docx + reportlab.
        Preserves document structure: inline tables, bold/italic, lists, images.
        """
        if not REPORTLAB_AVAILABLE:
            raise Exception(
                f"Required library reportlab not installed: {_reportlab_import_error}\n"
                "Please install: pip install reportlab"
            )
        if not DOCX_AVAILABLE:
            raise Exception(
                f"Required library python-docx not installed: {_docx_import_error}\n"
                "Please install: pip install python-docx"
            )

        try:
            # Import docx internals for document order iteration
            from docx.text.paragraph import Paragraph as DocxParagraph
            from docx.table import Table as DocxTable
            from docx.oxml.ns import qn
            import io

            # Read the DOCX file
            doc = Document(input_path)
            self.log.debug(f"Document loaded. Paragraphs: {len(doc.paragraphs)}, Tables: {len(doc.tables)}")

            # Register Unicode fonts
            self._register_unicode_fonts()

            # Create PDF in memory
            pdf_buffer = BytesIO()
            pdf_doc = SimpleDocTemplate(
                pdf_buffer,
                pagesize=letter,
                rightMargin=36,
                leftMargin=36,
                topMargin=36,
                bottomMargin=36
            )

            # Get default styles
            styles = getSampleStyleSheet()

            # Determine which font to use (prefer Unicode fonts if registered)
            font_name = 'UnicodeSans' if 'UnicodeSans' in pdfmetrics.getRegisteredFontNames() else 'Helvetica'
            font_name_bold = 'UnicodeSansBold' if 'UnicodeSansBold' in pdfmetrics.getRegisteredFontNames() else 'Helvetica-Bold'

            # Create custom styles matching document structure
            normal_style = ParagraphStyle(
                'CustomNormal',
                parent=styles['Normal'],
                fontName=font_name,
                fontSize=10,
                leading=12,
                spaceAfter=6
            )

            # List styles with proper indentation
            list_bullet_style = ParagraphStyle(
                'CustomListBullet',
                parent=styles['Normal'],
                fontName=font_name,
                fontSize=10,
                leading=12,
                spaceAfter=3,
                leftIndent=18,
                bulletIndent=6
            )

            list_bullet_2_style = ParagraphStyle(
                'CustomListBullet2',
                parent=styles['Normal'],
                fontName=font_name,
                fontSize=10,
                leading=12,
                spaceAfter=3,
                leftIndent=36,
                bulletIndent=24
            )

            list_number_style = ParagraphStyle(
                'CustomListNumber',
                parent=styles['Normal'],
                fontName=font_name,
                fontSize=10,
                leading=12,
                spaceAfter=3,
                leftIndent=18,
                bulletIndent=6
            )

            list_number_2_style = ParagraphStyle(
                'CustomListNumber2',
                parent=styles['Normal'],
                fontName=font_name,
                fontSize=10,
                leading=12,
                spaceAfter=3,
                leftIndent=36,
                bulletIndent=24
            )

            heading1_style = ParagraphStyle(
                'CustomHeading1',
                parent=styles['Heading1'],
                fontName=font_name_bold,
                fontSize=14,
                leading=18,
                spaceAfter=6,
                spaceBefore=10,
                textColor=colors.HexColor('#365F91')
            )

            heading2_style = ParagraphStyle(
                'CustomHeading2',
                parent=styles['Heading2'],
                fontName=font_name_bold,
                fontSize=12,
                leading=15,
                spaceAfter=4,
                spaceBefore=8,
                textColor=colors.HexColor('#4F81BD')
            )

            heading3_style = ParagraphStyle(
                'CustomHeading3',
                parent=styles['Heading3'],
                fontName=font_name_bold,
                fontSize=11,
                leading=14,
                spaceAfter=3,
                spaceBefore=6,
                textColor=colors.HexColor('#4F81BD')
            )

            # Code/monospace style for inline code and filenames
            font_name_mono = 'Courier'
            code_style = ParagraphStyle(
                'CustomCode',
                parent=styles['Normal'],
                fontName=font_name_mono,
                fontSize=9,
                leading=11,
                spaceAfter=4,
                backColor=colors.HexColor('#f5f5f5'),
                leftIndent=6,
                rightIndent=6
            )

            # Build the story (content) - iterate body elements in document order
            story = []

            def get_list_info(para):
                """Get list type and level from paragraph style and indentation."""
                style_name = para.style.name if para.style else ''

                # Determine list type from style name
                list_type = None
                if 'List Number' in style_name:
                    list_type = 'number'
                elif 'List Bullet' in style_name:
                    list_type = 'bullet'
                elif 'List' in style_name:
                    list_type = 'bullet'

                if list_type is None:
                    # Check numPr for lists without explicit style
                    try:
                        if para._element.pPr is not None and para._element.pPr.numPr is not None:
                            list_type = 'bullet'
                    except AttributeError:
                        pass

                if list_type is None:
                    return (None, 0)

                # Determine level from style name (List Number 2, List Bullet 2)
                if '2' in style_name or '3' in style_name:
                    level = 1
                else:
                    # Check leftIndent for nesting level
                    level = 0
                    try:
                        pPr = para._element.pPr
                        if pPr is not None:
                            ind = pPr.find(qn('w:ind'))
                            if ind is not None:
                                left_val = ind.get(qn('w:left'))
                                if left_val:
                                    left_indent = int(left_val)
                                    # 720 twips = level 0, 1440+ = level 1+
                                    if left_indent > 720:
                                        level = 1
                    except (AttributeError, ValueError):
                        pass

                return (list_type, level)

            # Track numbering for ordered lists
            number_counters = {0: 0, 1: 0, 2: 0}
            last_list_level = -1

            def is_code_run(run):
                """Check if a run has code/monospace styling."""
                try:
                    # Check style name for code indicators
                    if run.style and run.style.name:
                        style_name = run.style.name.lower()
                        if any(kw in style_name for kw in ['code', 'verbatim', 'mono', 'console']):
                            return True

                    # Check font name for monospace fonts
                    if run.font and run.font.name:
                        font_name_lower = run.font.name.lower()
                        if any(kw in font_name_lower for kw in ['courier', 'consolas', 'mono', 'code']):
                            return True

                    # Check XML for rFonts element with monospace font
                    rPr = run._element.rPr
                    if rPr is not None:
                        rFonts = rPr.find(qn('w:rFonts'))
                        if rFonts is not None:
                            ascii_font = rFonts.get(qn('w:ascii'))
                            if ascii_font:
                                font_lower = ascii_font.lower()
                                if any(kw in font_lower for kw in ['courier', 'consolas', 'mono', 'code']):
                                    return True
                except (AttributeError, TypeError):
                    pass
                return False

            def is_horizontal_rule(para):
                """Check if paragraph represents a horizontal divider line."""
                try:
                    # Check for paragraph border (horizontal line)
                    pPr = para._element.pPr
                    if pPr is not None:
                        pBdr = pPr.find(qn('w:pBdr'))
                        if pBdr is not None:
                            # Check for bottom border which creates a horizontal line
                            bottom = pBdr.find(qn('w:bottom'))
                            top = pBdr.find(qn('w:top'))
                            if bottom is not None or top is not None:
                                # If paragraph is empty or just whitespace, it's a divider
                                if not para.text.strip():
                                    return True
                except (AttributeError, TypeError):
                    pass
                return False

            def process_paragraph(para):
                """Process a single paragraph and return reportlab element(s)."""
                nonlocal last_list_level

                # Check for horizontal rule/divider first
                if is_horizontal_rule(para):
                    from reportlab.platypus import HRFlowable
                    return [HRFlowable(width="100%", thickness=0.5, color=colors.grey, spaceBefore=3, spaceAfter=6)]

                text = para.text.strip()
                if not text:
                    # Empty paragraph - use line-height spacing for proper separation
                    return Spacer(1, 0.15 * inch)

                # Escape XML special characters for base text and convert newlines
                text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                text = text.replace('\n', '<br/>')

                # Process all runs with their formatting
                def format_run(run):
                    """Format a single run with all its styling."""
                    run_text = run.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    if not run_text:
                        return run_text

                    # Convert newlines to HTML line breaks for reportlab
                    run_text = run_text.replace('\n', '<br/>')

                    result = run_text

                    # Check for code/monospace first (takes precedence)
                    if is_code_run(run):
                        result = f'<font face="Courier" size="9">{result}</font>'
                        return result

                    # Apply formatting tags (can be combined)
                    if run.bold:
                        result = f'<b>{result}</b>'
                    if run.italic:
                        result = f'<i>{result}</i>'
                    if run.underline:
                        result = f'<u>{result}</u>'
                    if run.font.strike:
                        result = f'<strike>{result}</strike>'
                    if run.font.subscript:
                        result = f'<sub>{result}</sub>'
                    if run.font.superscript:
                        result = f'<super>{result}</super>'

                    # Check for text color
                    try:
                        if run.font.color and run.font.color.rgb:
                            color_hex = str(run.font.color.rgb)
                            if color_hex and color_hex != 'None' and len(color_hex) == 6:
                                result = f'<font color="#{color_hex}">{result}</font>'
                    except (AttributeError, TypeError):
                        pass

                    return result

                # Check if any run has formatting that needs processing
                has_formatting = False
                for run in para.runs:
                    if run.text.strip():
                        if (run.bold or run.italic or run.underline or
                            run.font.strike or run.font.subscript or run.font.superscript or
                            is_code_run(run)):
                            has_formatting = True
                            break
                        try:
                            if run.font.color and run.font.color.rgb:
                                has_formatting = True
                                break
                        except (AttributeError, TypeError):
                            pass

                if has_formatting:
                    formatted_parts = [format_run(run) for run in para.runs]
                    text = ''.join(formatted_parts)

                # Detect heading styles
                style_name = para.style.name if para.style else ''
                if style_name.startswith('Heading 1'):
                    last_list_level = -1
                    return Paragraph(text, heading1_style)
                elif style_name.startswith('Heading 2'):
                    last_list_level = -1
                    return Paragraph(text, heading2_style)
                elif style_name.startswith('Heading 3') or style_name.startswith('Heading'):
                    last_list_level = -1
                    return Paragraph(text, heading3_style)

                # Check for list items
                list_type, level = get_list_info(para)

                if list_type == 'number':
                    # Reset lower levels when moving up, increment current level
                    if level <= last_list_level:
                        for l in range(level + 1, 3):
                            number_counters[l] = 0
                    number_counters[level] += 1
                    last_list_level = level

                    prefix = f"{number_counters[level]}. "
                    style = list_number_2_style if level > 0 else list_number_style
                    return Paragraph(f'{prefix}{text}', style)

                elif list_type == 'bullet':
                    last_list_level = level
                    style = list_bullet_2_style if level > 0 else list_bullet_style
                    return Paragraph(f'• {text}', style)

                else:
                    # Reset counters when not in list
                    last_list_level = -1
                    for l in number_counters:
                        number_counters[l] = 0
                    return Paragraph(text, normal_style)

            def process_table(tbl):
                """Process a single table and return reportlab elements."""
                table_data = []
                for row in tbl.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = cell.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        row_data.append(cell_text)
                    table_data.append(row_data)

                if not table_data:
                    return []

                t = Table(table_data, hAlign='LEFT')
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#dbe5f1')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor('#365F91')),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), font_name_bold),
                    ('FONTNAME', (0, 1), (-1, -1), font_name),
                    ('FONTSIZE', (0, 0), (-1, -1), 9),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
                    ('TOPPADDING', (0, 0), (-1, -1), 4),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#cccccc'))
                ]))
                return [t, Spacer(1, 0.15 * inch)]

            def process_image(drawing_element, doc):
                """Extract image from drawing element and return reportlab Image."""
                if not PIL_AVAILABLE:
                    return None
                try:
                    # Navigate to blip element containing image reference
                    blip = drawing_element.find('.//' + qn('a:blip'))
                    if blip is None:
                        return None

                    # Get the relationship ID
                    rId = blip.get(qn('r:embed'))
                    if not rId:
                        return None

                    # Get image data from document parts
                    image_part = doc.part.related_parts.get(rId)
                    if not image_part:
                        return None

                    # Create reportlab Image from bytes
                    img_buffer = io.BytesIO(image_part.blob)
                    img = RLImage(img_buffer)

                    # Scale image to fit page width (max 7 inches)
                    max_width = 7 * inch
                    if img.drawWidth > max_width:
                        scale = max_width / img.drawWidth
                        img.drawWidth = max_width
                        img.drawHeight = img.drawHeight * scale

                    # Left-align image
                    img.hAlign = 'LEFT'

                    return img
                except Exception as e:
                    self.log.debug(f"Error extracting image: {e}")
                    return None

            def add_to_story(result):
                """Add paragraph result to story, handling lists or single elements."""
                if isinstance(result, list):
                    story.extend(result)
                else:
                    story.append(result)

            # Iterate through body elements in document order (preserves table position)
            para_count = 0
            table_count = 0
            for element in doc.element.body:
                if element.tag == qn('w:p'):  # Paragraph
                    para = DocxParagraph(element, doc)
                    para_count += 1

                    # Check for drawings (images) in paragraph
                    drawings = element.findall('.//' + qn('w:drawing'))
                    if drawings:
                        # Process paragraph text first (if any)
                        if para.text.strip():
                            add_to_story(process_paragraph(para))
                        # Then add images
                        for drawing in drawings:
                            img = process_image(drawing, doc)
                            if img:
                                story.append(img)
                                story.append(Spacer(1, 0.1 * inch))
                    else:
                        add_to_story(process_paragraph(para))

                elif element.tag == qn('w:tbl'):  # Table
                    tbl = DocxTable(element, doc)
                    table_count += 1
                    story.extend(process_table(tbl))

            self.log.debug(f"Processed {para_count} paragraphs, {table_count} tables in document order")

            # Build the PDF
            if not story:
                self.log.warning("Story is empty! Creating a placeholder PDF")
                story.append(Paragraph("Document appears to be empty or contains no readable content.", normal_style))

            pdf_doc.build(story)

            # Get the PDF bytes
            pdf_bytes = pdf_buffer.getvalue()
            pdf_buffer.close()

            return pdf_bytes

        except Exception as e:
            raise Exception(f"DOCX to PDF conversion error: {str(e)}")

    def _convert_pptx_to_pdf(self, input_path: str) -> bytes:
        """
        Convert PPTX presentation to PDF using python-pptx + reportlab.
        Renders each slide as a PDF page with text, shapes, and images.
        """
        if not REPORTLAB_AVAILABLE:
            raise Exception(
                f"Required library reportlab not installed: {_reportlab_import_error}\n"
                "Please install: pip install reportlab"
            )
        if not PPTX_AVAILABLE:
            raise Exception(
                f"Required library python-pptx not installed: {_pptx_import_error}\n"
                "Please install: pip install python-pptx"
            )

        try:
            # Load the presentation
            prs = Presentation(input_path)
            self.log.debug(f"Presentation loaded. Slides: {len(prs.slides)}")

            # Get slide dimensions (default is 10x7.5 inches for 4:3)
            slide_width_emu = prs.slide_width
            slide_height_emu = prs.slide_height

            # Convert EMU to points (1 inch = 914400 EMU, 1 inch = 72 points)
            slide_width_pt = slide_width_emu / 914400 * 72
            slide_height_pt = slide_height_emu / 914400 * 72

            self.log.debug(f"Slide size: {slide_width_pt:.1f} x {slide_height_pt:.1f} points")

            # Register Unicode fonts (same as DOCX handler)
            self._register_unicode_fonts()

            # Determine which font to use
            font_name = 'UnicodeSans' if 'UnicodeSans' in pdfmetrics.getRegisteredFontNames() else 'Helvetica'
            font_name_bold = 'UnicodeSansBold' if 'UnicodeSansBold' in pdfmetrics.getRegisteredFontNames() else 'Helvetica-Bold'

            # Create PDF in memory using canvas for precise positioning
            pdf_buffer = BytesIO()
            c = canvas.Canvas(pdf_buffer, pagesize=(slide_width_pt, slide_height_pt))

            for slide_idx, slide in enumerate(prs.slides):
                self.log.debug(f"Processing slide {slide_idx + 1}")

                # Draw slide background (white by default)
                c.setFillColor(colors.white)
                c.rect(0, 0, slide_width_pt, slide_height_pt, fill=1, stroke=0)

                # Try to get slide background color
                try:
                    if slide.background.fill.solid():
                        bg_color = slide.background.fill.fore_color.rgb
                        if bg_color:
                            c.setFillColor(colors.HexColor(f'#{bg_color}'))
                            c.rect(0, 0, slide_width_pt, slide_height_pt, fill=1, stroke=0)
                except Exception:
                    pass  # Use default white background

                # Process shapes on the slide
                for shape in slide.shapes:
                    try:
                        self._render_shape_to_canvas(
                            c, shape, slide_width_emu, slide_height_emu,
                            slide_width_pt, slide_height_pt,
                            font_name, font_name_bold
                        )
                    except Exception as shape_error:
                        self.log.debug(f"Error rendering shape: {shape_error}")
                        continue

                # Add slide number at bottom
                c.setFont(font_name, 8)
                c.setFillColor(colors.grey)
                c.drawCentredString(slide_width_pt / 2, 15, f"Slide {slide_idx + 1}")

                # Move to next page
                c.showPage()

            # Save the PDF
            c.save()
            pdf_bytes = pdf_buffer.getvalue()
            pdf_buffer.close()

            self.log.info(f"PPTX conversion complete: {len(prs.slides)} slides")
            return pdf_bytes

        except Exception as e:
            raise Exception(f"PPTX to PDF conversion error: {str(e)}")

    def _register_unicode_fonts(self):
        """Register Unicode-supporting fonts from system paths."""
        from reportlab.pdfbase.pdfmetrics import registerFontFamily

        # Define font sets with normal, bold, italic, bolditalic variants
        font_sets = [
            # DejaVu fonts (most common, excellent Unicode support)
            {
                'normal': '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
                'bold': '/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf',
                'italic': None,  # Not available in this set
                'boldItalic': None,
            },
            # Liberation fonts (alternative)
            {
                'normal': '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
                'bold': '/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf',
                'italic': '/usr/share/fonts/truetype/liberation/LiberationSans-Italic.ttf',
                'boldItalic': '/usr/share/fonts/truetype/liberation/LiberationSans-BoldItalic.ttf',
            },
            # FreeSans (GNU FreeFont)
            {
                'normal': '/usr/share/fonts/truetype/freefont/FreeSans.ttf',
                'bold': '/usr/share/fonts/truetype/freefont/FreeSansBold.ttf',
                'italic': '/usr/share/fonts/truetype/freefont/FreeSansOblique.ttf',
                'boldItalic': '/usr/share/fonts/truetype/freefont/FreeSansBoldOblique.ttf',
            },
        ]

        font_names = {
            'normal': 'UnicodeSans',
            'bold': 'UnicodeSansBold',
            'italic': 'UnicodeSansItalic',
            'boldItalic': 'UnicodeSansBoldItalic',
        }

        registered_fonts = set()

        # Try each font set until we find one with at least normal and bold
        for font_set in font_sets:
            if 'UnicodeSans' in registered_fonts:
                break  # Already have fonts registered

            # Check if at least normal exists
            if font_set['normal'] and os.path.exists(font_set['normal']):
                for variant, path in font_set.items():
                    if path and os.path.exists(path):
                        font_name = font_names[variant]
                        if font_name not in registered_fonts:
                            try:
                                pdfmetrics.registerFont(TTFont(font_name, path))
                                registered_fonts.add(font_name)
                                self.log.info(f"Registered font {font_name} from {path}")
                            except Exception as font_error:
                                self.log.debug(f"Failed to register {font_name} from {path}: {font_error}")

        # Register font family to enable <b> and <i> tags in Paragraph
        if 'UnicodeSans' in registered_fonts:
            try:
                # Use Helvetica-Oblique as fallback for italic if no Unicode italic available
                # (Helvetica is a built-in PDF font, always available)
                italic_font = 'UnicodeSansItalic' if 'UnicodeSansItalic' in registered_fonts else 'Helvetica-Oblique'
                bold_italic_font = 'UnicodeSansBoldItalic' if 'UnicodeSansBoldItalic' in registered_fonts else 'Helvetica-BoldOblique'

                registerFontFamily(
                    'UnicodeSans',
                    normal='UnicodeSans',
                    bold='UnicodeSansBold' if 'UnicodeSansBold' in registered_fonts else 'UnicodeSans',
                    italic=italic_font,
                    boldItalic=bold_italic_font
                )
                self.log.info(f"Registered UnicodeSans font family (italic={italic_font})")
            except Exception as e:
                self.log.debug(f"Failed to register font family: {e}")

        if not registered_fonts:
            self.log.warning("No Unicode fonts found. International characters may not display correctly.")

    def _render_shape_to_canvas(self, c, shape, slide_width_emu, slide_height_emu,
                                 slide_width_pt, slide_height_pt, font_name, font_name_bold):
        """Render a single shape to the PDF canvas."""
        # Get shape position and size in points
        # Note: PDF coordinates start from bottom-left, PPTX from top-left
        left_pt = shape.left / 914400 * 72
        top_pt = shape.top / 914400 * 72
        width_pt = shape.width / 914400 * 72
        height_pt = shape.height / 914400 * 72

        # Convert to PDF coordinates (flip Y axis)
        pdf_y = slide_height_pt - top_pt - height_pt

        # Handle different shape types
        if shape.has_text_frame:
            self._render_text_frame(c, shape.text_frame, left_pt, pdf_y, width_pt, height_pt,
                                    slide_height_pt, font_name, font_name_bold)

        # Handle images
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            self._render_picture(c, shape, left_pt, pdf_y, width_pt, height_pt)

        # Handle tables
        if shape.has_table:
            self._render_table(c, shape.table, left_pt, pdf_y, width_pt, height_pt,
                              font_name, font_name_bold)

    def _render_text_frame(self, c, text_frame, x, y, width, height, slide_height_pt,
                           font_name, font_name_bold):
        """Render a text frame to the canvas."""
        current_y = y + height - 5  # Start from top of text box

        for paragraph in text_frame.paragraphs:
            text = paragraph.text.strip()
            if not text:
                current_y -= 12  # Empty paragraph spacing
                continue

            # Determine font size and style
            font_size = 12  # Default
            current_font = font_name

            try:
                if paragraph.runs:
                    run = paragraph.runs[0]
                    if run.font.size:
                        font_size = run.font.size.pt
                    if run.font.bold:
                        current_font = font_name_bold
            except Exception:
                pass

            # Clamp font size to reasonable range
            font_size = max(6, min(72, font_size))

            # Set font color
            try:
                if paragraph.runs and paragraph.runs[0].font.color.rgb:
                    rgb = paragraph.runs[0].font.color.rgb
                    c.setFillColor(colors.HexColor(f'#{rgb}'))
                else:
                    c.setFillColor(colors.black)
            except Exception:
                c.setFillColor(colors.black)

            c.setFont(current_font, font_size)

            # Handle text alignment
            try:
                from pptx.enum.text import PP_ALIGN
                if paragraph.alignment == PP_ALIGN.CENTER:
                    c.drawCentredString(x + width / 2, current_y, text)
                elif paragraph.alignment == PP_ALIGN.RIGHT:
                    c.drawRightString(x + width, current_y, text)
                else:
                    c.drawString(x + 5, current_y, text)
            except Exception:
                c.drawString(x + 5, current_y, text)

            # Move to next line
            current_y -= font_size * 1.2

    def _render_picture(self, c, shape, x, y, width, height):
        """Render a picture shape to the canvas."""
        if not PIL_AVAILABLE:
            self.log.debug("PIL not available, skipping image")
            return

        try:
            # Get image blob from shape
            image = shape.image
            image_bytes = image.blob

            # Create PIL image and save to buffer
            pil_image = PILImage.open(BytesIO(image_bytes))

            # Convert to RGB if necessary (for PDF compatibility)
            if pil_image.mode in ('RGBA', 'P'):
                pil_image = pil_image.convert('RGB')

            # Save to buffer for reportlab
            img_buffer = BytesIO()
            pil_image.save(img_buffer, format='PNG')
            img_buffer.seek(0)

            # Draw image on canvas using ImageReader
            from reportlab.lib.utils import ImageReader
            img_reader = ImageReader(img_buffer)
            c.drawImage(img_reader, x, y, width=width, height=height,
                        preserveAspectRatio=True, mask='auto')
        except Exception as e:
            self.log.debug(f"Error rendering image: {e}")
            # Draw placeholder rectangle
            c.setStrokeColor(colors.grey)
            c.setFillColor(colors.lightgrey)
            c.rect(x, y, width, height, fill=1, stroke=1)
            c.setFillColor(colors.grey)
            c.setFont('Helvetica', 8)
            c.drawCentredString(x + width/2, y + height/2, "[Image]")

    def _render_table(self, c, table, x, y, width, height, font_name, font_name_bold):
        """Render a table to the canvas."""
        if not table.rows:
            return

        num_rows = len(table.rows)
        num_cols = len(table.columns)

        cell_width = width / num_cols
        cell_height = height / num_rows

        current_y = y + height  # Start from top

        for row_idx, row in enumerate(table.rows):
            current_x = x
            current_y -= cell_height

            for col_idx, cell in enumerate(row.cells):
                # Draw cell border
                c.setStrokeColor(colors.black)
                c.rect(current_x, current_y, cell_width, cell_height, fill=0, stroke=1)

                # Draw cell text
                text = cell.text.strip()
                if text:
                    # Use bold for header row
                    if row_idx == 0:
                        c.setFont(font_name_bold, 9)
                        c.setFillColor(colors.HexColor('#333333'))
                    else:
                        c.setFont(font_name, 9)
                        c.setFillColor(colors.black)

                    # Truncate text if too long
                    max_chars = int(cell_width / 5)
                    if len(text) > max_chars:
                        text = text[:max_chars-2] + '..'

                    c.drawString(current_x + 3, current_y + cell_height/2 - 3, text)

                current_x += cell_width


def setup_handlers(web_app):
    host_pattern = ".*$"

    base_url = web_app.settings["base_url"]

    # Document converter endpoint
    converter_pattern = url_path_join(base_url, "jupyterlab-doc-reader-extension", "convert")
    handlers = [(converter_pattern, DocumentConverterHandler)]

    web_app.add_handlers(host_pattern, handlers)
